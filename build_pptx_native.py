#!/usr/bin/env python3
"""
Native (browser-free) 16:9 PPTX rebuild of the LVT Claude Code KASE deck.
Reads Decks/lvt_claude_code/slides_data.json (the exact SLIDE_DATA content)
and rebuilds each slide with real text/shapes in the KASE dark-blue palette.
Editable in PowerPoint; on-brand; no headless browser required.
"""
import json, os, re
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DECK = "Decks/lvt_claude_code"
DATA = json.load(open(os.path.join(DECK, "slides_data.json")))
OUT  = "LVT_AI_Assisted_Dev_Bootcamp_ClaudeCode.pptx"
LOGO = "deck/Images/lvt-logo.png"

# ---- palette (KASE Edelweiss-aligned dark-blue) ----
DARK   = RGBColor(0x26,0x25,0x33)
CORAL  = RGBColor(0xf8,0x6b,0x3c)
INDIGO = RGBColor(0x61,0x67,0xfa)
LIGHT  = RGBColor(0xf7,0xf9,0xfc)
PANEL  = RGBColor(0xee,0xf1,0xf7)
BORDER = RGBColor(0xdc,0xe2,0xee)
SLATE  = RGBColor(0x47,0x55,0x69)
MUTED  = RGBColor(0x64,0x74,0x8b)
NEARW  = RGBColor(0xf4,0xf6,0xfb)
WHITE  = RGBColor(0xff,0xff,0xff)
SERIF  = "DM Serif Display"
SANS   = "DM Sans"

SW, SH = 12192000, 6858000  # 13.333 x 7.5 in
IN = 914400
prs = Presentation(); prs.slide_width = Emu(SW); prs.slide_height = Emu(SH)
BLANK = prs.slide_layouts[6]

def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, Emu(SW), Emu(SH))
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s

def strip_fmt(t):
    if t is None: return ""
    t = str(t)
    t = re.sub(r"\[#?[a-zA-Z0-9]+:([^\]]*)\]", r"\1", t)  # [navy:x] / [#hex:x]
    t = re.sub(r"\[([^\]]*)\]", r"\1", t)                   # [x]
    t = t.replace("**","").replace("*","")
    return t

def box(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Emu(int(x*IN)), Emu(int(y*IN)), Emu(int(w*IN)), Emu(int(h*IN)))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tb, tf

def setpara(p, text, size, color, bold=False, font=SANS, italic=False, align=PP_ALIGN.LEFT, space_after=4, line=None):
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
    if line: p.line_spacing = line
    r = p.add_run(); r.text = strip_fmt(text); f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font
    f.color.rgb = color
    return p

def addtext(s, x, y, w, h, text, size, color, **kw):
    anchor = kw.pop("anchor", MSO_ANCHOR.TOP)
    tb, tf = box(s, x, y, w, h, anchor)
    setpara(tf.paragraphs[0], text, size, color, **kw)
    return tb

def rect(s, x, y, w, h, fill, line=None, line_w=1.0, rounded=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Emu(int(x*IN)), Emu(int(y*IN)), Emu(int(w*IN)), Emu(int(h*IN)))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def eyebrow(s, x, y, w, text, color=CORAL):
    addtext(s, x, y, w, 0.3, (text or "").upper(), 12.5, color, bold=True, font=SANS)

def footer(s, dark=False):
    c = MUTED if not dark else RGBColor(0x9a,0x9f,0xb5)
    addtext(s, 0.55, 7.06, 6, 0.3, "Acceler  |  Proprietary and confidential", 9, c, font=SANS)

def lines_from(content):
    """split a multi-line field into (bullets, plain) honoring '--' bullet markers."""
    out=[]
    for ln in str(content).split("\n"):
        ln=ln.strip()
        if not ln: continue
        if ln.startswith("--"): out.append(("bullet", ln[2:].strip()))
        else: out.append(("text", ln))
    return out

# ───────────────────────── templates ─────────────────────────
def t_cover(d):
    s = slide(DARK)
    rect(s, 0, 0, 0.16, 7.5, CORAL)
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Emu(int(0.75*IN)), Emu(int(0.7*IN)), height=Emu(int(0.62*IN)))
    addtext(s, 0.75, 1.7, 11, 0.4, d.get("tagline_prefix",""), 15, RGBColor(0xc7,0xcc,0xe0), font=SANS)
    tb,tf = box(s, 0.72, 2.15, 11.3, 2.2)
    for i,ln in enumerate(strip_fmt(d.get("headline","")).split("\n")):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        setpara(p, ln, 46, NEARW, font=SERIF, line=1.02, space_after=0)
    addtext(s, 0.75, 4.35, 11, 0.4, d.get("tagline_highlight",""), 15, CORAL, bold=True, font=SANS)
    # stats strip
    stats = [("👩‍💻","Audience","Software Developers"),("🤖","Primary Tool","Claude Code"),
             ("📅","Duration","3 Days · Live"),("⏱","Hours","18 Hrs Hands-On")]
    n=len(stats); gap=0.18; x0=0.75; tw=11.0; cw=(tw-gap*(n-1))/n
    for i,(ic,lb,vl) in enumerate(stats):
        x=x0+i*(cw+gap)
        rect(s, x, 5.25, cw, 1.15, RGBColor(0x32,0x31,0x40), rounded=True)
        addtext(s, x+0.12, 5.4, cw-0.24, 0.3, lb.upper(), 9.5, CORAL, bold=True, font=SANS)
        addtext(s, x+0.12, 5.72, cw-0.24, 0.55, vl, 14, NEARW, bold=True, font=SANS, anchor=MSO_ANCHOR.TOP)
    addtext(s, 0.75, 6.62, 6, 0.3, d.get("date",""), 12, RGBColor(0x9a,0x9f,0xb5), font=SANS)
    addtext(s, 6, 6.62, 6, 0.3, "BY UnO · UPWARD AND ONWARD", 10, RGBColor(0x9a,0x9f,0xb5), font=SANS, align=PP_ALIGN.RIGHT)

def t_divider(d):
    s = slide(DARK)
    rect(s, 0, 0, 0.16, 7.5, CORAL)
    eyebrow(s, 0.85, 2.3, 6, d.get("section_label",""))
    tb,tf = box(s, 0.82, 2.7, 10.5, 2.0)
    for i,ln in enumerate(strip_fmt(d.get("headline","")).split("\n")):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        setpara(p, ln, 44, NEARW, font=SERIF, line=1.03, space_after=0)
    if d.get("subtext"):
        addtext(s, 0.85, 4.7, 10.0, 1.2, d["subtext"], 15, RGBColor(0xc7,0xcc,0xe0), font=SANS, line=1.3)
    footer(s, dark=True)

def header(s, d):
    eyebrow(s, 0.55, 0.5, 11, d.get("eyebrow",""))
    tb,tf = box(s, 0.52, 0.82, 12.0, 1.0)
    setpara(tf.paragraphs[0], d.get("headline",""), 30, DARK, font=SERIF, line=1.02)

def t_acceler_intro(d):
    s = slide(LIGHT); header(s, d)
    stats=d.get("stats",[])
    cols=3; gap=0.25; x0=0.55; tw=12.2; cw=(tw-gap*(cols-1))/cols; ch=1.18
    for i,st in enumerate(stats[:6]):
        r=i//cols; c=i%cols; x=x0+c*(cw+gap); y=1.95+r*(ch+0.22)
        rect(s, x, y, cw, ch, PANEL, rounded=True)
        parts=str(st.get("content","")).split("\n")
        num=parts[0]; lab=parts[1][2:].strip() if len(parts)>1 else ""
        addtext(s, x+0.18, y+0.16, cw-0.36, 0.5, num, 26, INDIGO, bold=True, font=SERIF)
        addtext(s, x+0.18, y+0.66, cw-0.36, 0.45, lab, 11, SLATE, font=SANS, line=1.05)
    addtext(s, 0.55, 6.55, 12, 0.3, "Delivered by the world's top AI practitioners — FAANG & top AI labs", 11, MUTED, italic=True, font=SANS)
    footer(s)

def t_approach(d):
    s = slide(LIGHT); header(s, d)
    pts=d.get("points",[]); cols=2; gap=0.3; x0=0.55; tw=12.2; cw=(tw-gap*(cols-1))/cols; ch=1.15
    for i,pt in enumerate(pts[:6]):
        r=i//cols; c=i%cols; x=x0+c*(cw+gap); y=1.95+r*(ch+0.16)
        rect(s, x, y, cw, ch, WHITE, line=BORDER, rounded=True)
        rect(s, x, y, 0.07, ch, CORAL)
        addtext(s, x+0.22, y+0.14, 0.5, 0.4, pt.get("num",""), 20, CORAL, bold=True, font=SERIF)
        addtext(s, x+0.7, y+0.15, cw-0.9, 0.4, pt.get("title",""), 14.5, DARK, bold=True, font=SANS)
        addtext(s, x+0.7, y+0.52, cw-0.9, 0.55, pt.get("desc",""), 11, SLATE, font=SANS, line=1.12)
    footer(s)

def t_overview(d):
    s = slide(LIGHT); header(s, d)
    rows=d.get("rows",[])
    y=1.9; x=0.55; w=12.2; lw=2.4
    rect(s, x, y, w, 6.55-y-0.0, WHITE, line=BORDER, rounded=True)
    yy=y+0.18
    for row in rows:
        label=row.get("label",""); typ=row.get("type"); content=row.get("content","")
        # estimate height
        if typ=="bullets":
            items=[c for k,c in lines_from(content)]
            h=0.26*len(items)+0.18
        elif typ=="pills":
            h=0.5
        else:
            h=0.24*max(1, (len(strip_fmt(content))//95)+1)+0.18
        addtext(s, x+0.22, yy, lw-0.1, h, label, 11.5, CORAL, bold=True, font=SANS)
        cx=x+lw+0.1; cw=w-lw-0.4
        if typ=="pills":
            toks=[t.strip() for t in str(content).split("·") if t.strip()]
            px=cx; py=yy
            for tok in toks:
                tw_=min(0.13*len(tok)+0.3, cw)
                if px+tw_>cx+cw: px=cx; py+=0.42
                rect(s, px, py, tw_, 0.32, PANEL, rounded=True)
                addtext(s, px+0.08, py+0.04, tw_-0.16, 0.26, tok, 9.5, SLATE, bold=True, font=SANS)
                px+=tw_+0.12
            h=max(h, (py-yy)+0.45)
        elif typ=="bullets":
            tb,tf=box(s, cx, yy, cw, h)
            for j,(k,c) in enumerate(lines_from(content)):
                p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
                setpara(p, "▸ "+c, 11, SLATE, font=SANS, line=1.1, space_after=2)
        else:
            addtext(s, cx, yy, cw, h, content, 11.5, SLATE, font=SANS, line=1.18)
        yy+=h+0.12
    footer(s)

def t_module_table(d):
    s = slide(LIGHT); header(s, d)
    rows=d.get("rows",[])
    cols=[("Day",2.9,"module"),("What developers leave with",3.6,"outcome"),("Parts (Claude Code, hands-on)",4.0,"tags"),("Hrs",0.9,"duration")]
    x0=0.55; y=1.95; total=sum(c[1] for c in cols); scale=12.2/total
    # header row
    cx=x0
    for name,w,_ in cols:
        ww=w*scale; rect(s, cx, y, ww, 0.42, DARK)
        addtext(s, cx+0.1, y+0.08, ww-0.2, 0.3, name, 10.5, WHITE, bold=True, font=SANS); cx+=ww
    yy=y+0.42
    palette={"navy":INDIGO,"teal":RGBColor(0x0f,0x76,0x6e),"amber":RGBColor(0xb4,0x53,0x09)}
    for ri,row in enumerate(rows):
        # height by tags lines
        tlines=[l for l in str(row.get("tags","")).split("\n") if l.strip()]
        mlines=[l for l in str(row.get("module","")).split("\n") if l.strip()]
        h=max(0.28*len(tlines)+0.2, 1.05)
        bg=WHITE if ri%2==0 else PANEL
        rect(s, x0, yy, 12.2, h, bg, line=BORDER)
        accent=palette.get(row.get("color"),INDIGO)
        rect(s, x0, yy, 0.07, h, accent)
        cx=x0
        for name,w,key in cols:
            ww=w*scale
            val=row.get(key,"")
            tb,tf=box(s, cx+0.14, yy+0.1, ww-0.24, h-0.18, MSO_ANCHOR.TOP)
            vlines=[l for l in str(val).split("\n") if l.strip()] or [""]
            for j,vl in enumerate(vlines):
                p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
                bold = (key=="module" and j==0) or key=="duration"
                col = DARK if key in ("module","duration") else SLATE
                if key=="module" and j>0: col=MUTED
                setpara(p, vl, 10 if key!="module" else (11 if j==0 else 9.5), col, bold=bold, font=SANS, line=1.08, space_after=1)
            cx+=ww
        yy+=h
    footer(s)

def t_two_cards(d):
    s = slide(LIGHT); header(s, d)
    gap=0.4; x0=0.55; tw=12.2; cw=(tw-gap)/2; y=1.95; ch=3.7
    for i in (1,2):
        x=x0+(i-1)*(cw+gap)
        rect(s, x, y, cw, ch, WHITE, line=BORDER, rounded=True)
        rect(s, x, y, cw, 0.08, CORAL if i==1 else INDIGO)
        addtext(s, x+0.25, y+0.22, cw-0.5, 0.4, d.get(f"card_{i}_title",""), 16, DARK, bold=True, font=SANS)
        addtext(s, x+0.25, y+0.66, cw-0.5, 0.5, d.get(f"card_{i}_sub",""), 11.5, MUTED, italic=True, font=SANS, line=1.1)
        tb,tf=box(s, x+0.25, y+1.2, cw-0.5, ch-1.4)
        for j,(k,c) in enumerate(lines_from(d.get(f"card_{i}_content",""))):
            p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
            setpara(p, "▸ "+c, 11.5, SLATE, font=SANS, line=1.15, space_after=6)
    if d.get("takeaway"):
        rect(s, x0, y+ch+0.2, tw, 0.7, DARK, rounded=True)
        addtext(s, x0+0.25, y+ch+0.34, tw-0.5, 0.45, d["takeaway"], 12.5, NEARW, bold=True, font=SANS, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)

def t_use_case(d):
    s = slide(LIGHT); header(s, d)
    if d.get("section_label"):
        addtext(s, 0.55, 1.78, 12, 0.3, d["section_label"], 11, MUTED, italic=True, font=SANS)
    cases=d.get("cases",[]); y=2.2; x0=0.55; tw=12.2; ch=1.0
    for c in cases:
        txt=str(c.get("content",""));
        head=txt.split("\n")[0]; body=txt.split("\n")[1] if "\n" in txt else ""
        title=head.split("|")[0].strip(); sub=head.split("|")[1].strip() if "|" in head else ""
        rect(s, x0, y, tw, ch, WHITE, line=BORDER, rounded=True)
        rect(s, x0, y, 0.07, ch, INDIGO)
        addtext(s, x0+0.25, y+0.13, tw-0.5, 0.35, title+("  —  "+sub if sub else ""), 13.5, DARK, bold=True, font=SANS)
        addtext(s, x0+0.25, y+0.5, tw-0.5, 0.45, body, 11, SLATE, font=SANS, line=1.12)
        y+=ch+0.16
    footer(s)

def t_mentor(d):
    s = slide(LIGHT)
    pw=4.0
    rect(s, 0, 0, pw, 7.5, DARK)
    # avatar
    init=d.get("avatar_initials","?")
    av=s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int((pw/2-0.85)*IN)), Emu(int(0.5*IN)), Emu(int(1.7*IN)), Emu(int(1.7*IN)))
    av.fill.solid(); av.fill.fore_color.rgb=CORAL; av.line.color.rgb=CORAL; av.shadow.inherit=False
    av.text_frame.text=init
    p=av.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    p.runs[0].font.size=Pt(40); p.runs[0].font.bold=True; p.runs[0].font.name=SERIF; p.runs[0].font.color.rgb=WHITE
    av.text_frame.word_wrap=True
    eyebrow(s, 0.45, 2.5, pw-0.8, d.get("spec_label","Technical Specializations"))
    tb,tf=box(s, 0.45, 2.9, pw-0.8, 3.0)
    for j,sp in enumerate(d.get("specializations",[])):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
        setpara(p, "◆ "+strip_fmt(sp), 11, NEARW, font=SANS, line=1.12, space_after=6)
    # companies
    addtext(s, 0.45, 6.2, pw-0.8, 0.9, "   ·   ".join(d.get("companies",[])), 10, RGBColor(0xb8,0xbd,0xd0), bold=True, font=SANS, line=1.2)
    # right
    rx=pw+0.55
    eyebrow(s, rx, 0.6, 7, d.get("eyebrow","Facilitator Profile"))
    addtext(s, rx, 1.0, 12.2-rx, 0.7, d.get("name",""), 36, DARK, font=SERIF)
    addtext(s, rx, 1.75, 12.2-rx, 0.7, d.get("role",""), 13, MUTED, bold=True, font=SANS, line=1.2)
    rect(s, rx, 2.55, 12.2-rx-0.3, 0.02, BORDER)
    eyebrow(s, rx, 2.75, 7, d.get("hl_label","Career Highlights"), color=INDIGO)
    tb,tf=box(s, rx, 3.2, 12.2-rx-0.3, 3.6)
    for j,h in enumerate(d.get("highlights",[])):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
        setpara(p, "▸ "+strip_fmt(h), 13, SLATE, font=SANS, line=1.2, space_after=9)
    addtext(s, rx, 7.06, 7, 0.3, d.get("footnote","*Indicative Facilitator Profile"), 9, MUTED, italic=True, font=SANS)

def t_closing(d):
    s = slide(DARK)
    rect(s, 0, 0, 0.16, 7.5, CORAL)
    addtext(s, 0.85, 2.4, 11, 1.0, d.get("tagline","Plug in and Accelerate"), 50, NEARW, font=SERIF, anchor=MSO_ANCHOR.TOP)
    addtext(s, 0.88, 3.6, 11, 0.4, "AI enabling engineering orgs for 100X productivity.", 14, CORAL, italic=True, font=SANS)
    y=4.4
    for c in d.get("contacts",[]):
        addtext(s, 0.88, y, 8, 0.35, c.get("name","")+"  ·  "+c.get("title",""), 14, NEARW, bold=True, font=SANS)
        addtext(s, 0.88, y+0.32, 8, 0.3, c.get("email",""), 12, RGBColor(0xc7,0xcc,0xe0), font=SANS)
        y+=0.85
    addtext(s, 0.85, 6.95, 11, 0.3, "BY UnO · UPWARD AND ONWARD", 10, RGBColor(0x9a,0x9f,0xb5), font=SANS)

TEMPLATES = {
    "cover_dark": t_cover, "section_divider_dark": t_divider, "acceler_intro": t_acceler_intro,
    "approach_points": t_approach, "solution_overview": t_overview, "module_table": t_module_table,
    "two_cards": t_two_cards, "use_case": t_use_case, "mentor_profile": t_mentor, "closing": t_closing,
}

for entry in DATA:
    d = entry["data"]; tpl = d.get("template")
    fn = TEMPLATES.get(tpl)
    if not fn:
        s = slide(LIGHT); addtext(s, 1,1,10,1, "Unmapped template: "+str(tpl), 20, CORAL)
        continue
    fn(d)

prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))
