#!/usr/bin/env python3
"""
Build a 16:9 PPTX from the rendered AI Builder HTML deck.
Each slide = one full-bleed 2x PNG (pixel-faithful to the live HTML design).
Slide order follows deck.json. No design changes — image-per-slide fidelity.
"""
import json, os, glob
from pptx import Presentation
from pptx.util import Emu

DECK_DIR = "Decks/lvt_claude_code"
SHOTS    = "/tmp/lvtshots"
OUT      = "LVT_AI_Assisted_Dev_Bootcamp_ClaudeCode.pptx"

# 16:9 widescreen — 13.333in x 7.5in in EMU
SW = 12192000
SH = 6858000

slides = json.load(open(os.path.join(DECK_DIR, "deck.json")))["slides"]

prs = Presentation()
prs.slide_width  = Emu(SW)
prs.slide_height = Emu(SH)
blank = prs.slide_layouts[6]  # blank layout

missing = []
for s in slides:
    name = s.replace(".html", "")
    img  = os.path.join(SHOTS, name + ".png")
    if not os.path.exists(img):
        missing.append(name); continue
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(img, Emu(0), Emu(0), width=Emu(SW), height=Emu(SH))

prs.save(OUT)
print("saved:", OUT)
print("slides written:", len(prs.slides.__iter__.__self__._sldIdLst))
print("missing:", missing or "none")
